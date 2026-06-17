"""
Generate a 10-slide overview deck for the Risk DQ Governance Copilot.

Portfolio / personal-reference deck — synthetic data, **no firm names, no PII**.
Run:  python make_ppt.py
Out:  Risk_DQ_Governance_Overview.pptx
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Theme ────────────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x0D, 0x2B, 0x55)
CYAN  = RGBColor(0x00, 0xAE, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCF, 0xE0, 0xF0)
MUTE  = RGBColor(0x8F, 0xA8, 0xC4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _accent_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0.0), Inches(0.0), Inches(0.18), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()


def _footer(slide, n):
    tf = _box(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"Risk DQ Governance Copilot · synthetic data · BCBS 239 + AI Governance        {n} / 10"
    r.font.size = Pt(10); r.font.color.rgb = MUTE


def title_slide():
    s = prs.slides.add_slide(BLANK)
    _bg(s); _accent_bar(s)
    tf = _box(s, Inches(0.9), Inches(2.2), Inches(11.6), Inches(2.6))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Risk Data Quality Governance Copilot"
    r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r = p2.add_run()
    r.text = "BCBS 239 risk-data quality, reporting & AI governance — a working demo"
    r.font.size = Pt(22); r.font.color.rgb = CYAN
    p3 = tf.add_paragraph()
    r = p3.add_run()
    r.text = "Synthetic data · Python · Streamlit · Claude · portfolio project"
    r.font.size = Pt(15); r.font.color.rgb = MUTE
    _footer(s, 1)


def content_slide(n, title, subtitle, bullets):
    """bullets: list of (text, level) — level 0 = main (cyan), 1 = sub (light)."""
    s = prs.slides.add_slide(BLANK)
    _bg(s); _accent_bar(s)
    # Title
    tf = _box(s, Inches(0.6), Inches(0.5), Inches(12.2), Inches(1.1))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        r = p2.add_run(); r.text = subtitle
        r.font.size = Pt(16); r.font.color.rgb = CYAN
    # Body
    body = _box(s, Inches(0.7), Inches(1.9), Inches(12.0), Inches(5.0))
    first = True
    for text, level in bullets:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.space_after = Pt(6)
        r = p.add_run()
        if level == 0:
            r.text = "▸  " + text
            r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE
            p.level = 0
        else:
            r.text = "•  " + text
            r.font.size = Pt(16); r.font.color.rgb = LIGHT
            p.level = 1
    _footer(s, n)


# ── Slides ───────────────────────────────────────────────────────────────────
title_slide()

content_slide(2, "The problem", "Why risk-data quality is a board-level capability", [
    ("After 2008, banks could not aggregate risk exposures fast or accurately", 0),
    ("\"What is our total exposure to this counterparty, right now?\" — many couldn't answer in time", 1),
    ("BCBS 239 (Basel, 2013): the standard for risk data aggregation & reporting (RDARR)", 0),
    ("Aggregate risk data accurately, completely, on time — especially under stress", 1),
    ("Still a top supervisory priority (ECB 2025–27); most banks not yet fully compliant", 0),
    ("And now AI raises both the opportunity and the risk", 0),
])

content_slide(3, "BCBS 239 in one slide", "14 principles, 4 groups", [
    ("Group I — Governance & infrastructure (P1–2)", 0),
    ("Ownership, data architecture, lineage, golden sources", 1),
    ("Group II — Risk data aggregation (P3–6)  ← the core", 0),
    ("Accuracy · Completeness · Timeliness · Adaptability", 1),
    ("Group III — Risk reporting (P7–11)", 0),
    ("Accuracy · comprehensiveness · clarity · frequency · distribution", 1),
    ("Group IV — Supervisory review (P12–14)", 0),
])

content_slide(4, "The solution", "Five components, one governed flow", [
    ("Detect → Score → Triage → Report → Govern", 0),
    ("1 · Synthetic risk dataset — credit/market/liquidity CDEs", 1),
    ("2 · DQ rules engine — mapped to BCBS 239 principles → RAG scorecard", 1),
    ("3 · AI triage agent — findings into governed DQGC issues", 1),
    ("4 · Reporting dashboard — Streamlit, decision-useful views", 1),
    ("5 · AI governance layer — 'AI is a model', controls & risk map", 1),
])

content_slide(5, "Data quality engine", "Components 1–2 · aggregation principles P3–P6", [
    ("606 synthetic exposures across three risk domains", 0),
    ("114 data-quality defects deliberately seeded and tagged to a principle", 1),
    ("Rules grouped by principle: Accuracy, Completeness, Timeliness, Adaptability", 0),
    ("RAG scorecard = 'evidence of compliance' for the DQGC and supervisors", 0),
    ("100% control coverage against the seeded defects (controls are provably correct)", 0),
])

content_slide(6, "AI exception triage", "Component 3 · governance P1 & P13", [
    ("120 raw findings collapsed into ~10 governed issues", 0),
    ("One issue per failed control population — not 120 tickets", 1),
    ("Each issue gets a DQGC-ready remediation narrative", 0),
    ("Root-cause hypothesis · risk impact · fix-at-source · owner · priority", 1),
    ("Claude drafts; graceful $0 templated fallback when offline", 0),
    ("\"AI drafts; the DQGC decides\" — human governance stays in control", 0),
])

content_slide(7, "Reporting dashboard", "Component 4 · reporting principles P7–P11", [
    ("Streamlit app that embodies the BCBS 239 reporting principles", 0),
    ("Live scorecard that recomputes from source — it can't drift (P7)", 1),
    ("Findings drill-down (P9) · DQGC issues · comprehensiveness across stripes (P8)", 1),
    ("Adaptability (P6): re-aggregate exposure on demand, no re-engineering", 0),
    ("One shareable, decision-useful view for senior management", 0),
])

content_slide(8, "AI governance layer", "Component 5 · 'AI is a model'", [
    ("AI touchpoint register — you can't govern what you can't see", 0),
    ("Honest control self-assessment (RAG) — gaps shown, not hidden", 0),
    ("Full 14-principle AI-risk → control map, tagged to the regime", 0),
    ("Provenance stamped on every AI output: model, version, 'Pending DQGC review'", 0),
    ("Aligned to EU AI Act · revised US Model Risk Management · BIS/ECB guidance", 0),
])

content_slide(9, "AI risk to BCBS 239", "The two-sided insight", [
    ("AI both threatens AND strengthens the principles", 0),
    ("Threat: hallucination → P3 · black-box transforms → P2 lineage · non-reconciling narrative → P7", 1),
    ("Enabler: DQ at scale · anomaly detection · natural-language querying (P6)", 1),
    ("AI risk is a lens over all 14 principles — not a 15th principle", 0),
    ("Anchors: EU AI Act (high-risk from Aug 2026) · US MRM (retiring SR 11-7) · BIS/ECB", 0),
])

content_slide(10, "Outcomes & roadmap", "What it demonstrates", [
    ("Four senior signals at once: domain depth · DQ governance · reporting modernisation · AI governance", 0),
    ("Honest framing: synthetic data, a thinking tool — not a production framework", 0),
    ("Roadmap: automated output reconciliation · drift monitoring · real-data integration", 0),
    ("Stack: Python · pandas · Streamlit · Claude (Anthropic)", 0),
])

OUT = "Risk_DQ_Governance_Overview.pptx"
prs.save(OUT)
print(f"Saved {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
